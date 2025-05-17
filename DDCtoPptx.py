#### Script for turning exports from DDC.R into pptx

import pandas as pd
import numpy as np
import glob
from pptx import Presentation
from pptx.util import Inches

# Directory path
directory_path = r'T:\DDC\R\csv\exports\DDC'

# Glob to find all csvs in the directory
distancereport = pd.read_csv(directory_path + r'\DistanceReport.csv')
leaderboard_for = pd.read_csv(directory_path + r'\leaderboard_for.csv')
leaderboard_tld = pd.read_csv(directory_path + r'\leaderboard_tld.csv')

df_leaderboard = pd.concat([leaderboard_for, leaderboard_tld])
print(df_leaderboard)

# Create pptx
prs = Presentation()

# Add a slide to the presentation
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 is a blank slide layout

# Create a table on the slide
graphic_frame = slide.shapes.add_table(rows=len(leaderboard_for)+1, cols=4, left=Inches(1), top=Inches(1), width=Inches(6), height=Inches(2))

# Create a table in the slide
table = graphic_frame.table

# Add df to the slide
for i, row in leaderboard_for.iterrows():
    table.cell(i+1, 0).text = str(row['bhid'])
    table.cell(i+1, 1).text = str(row['distance'])
    table.cell(i+1, 2).text = str(row['date'])
    table.cell(i+1, 3).text = str(row['rig_id'])

# Save presentation to a file
prs.save(directory_path + r'\leaderboard.pptx')